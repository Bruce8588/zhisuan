#!/usr/bin/env python3
"""创建专业的TD9策略路演PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import pandas as pd

# 颜色配置
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 102, 204)
GREEN = RGBColor(0, 153, 76)
RED = RGBColor(204, 0, 0)
GOLD = RGBColor(255, 153, 51)
WHITE = RGBColor(255, 255, 255)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_background(slide, color=DARK_BLUE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_box(slide, text, top=Inches(0.3), font_size=36):
    title_box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = WHITE
    return title_box

def add_content_box(slide, lines, top=Inches(1.5), font_size=22):
    content_box = slide.shapes.add_textbox(Inches(0.5), top, Inches(6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
    return content_box

def add_big_number(slide, number, label, left, top):
    num_box = slide.shapes.add_textbox(left, top, Inches(3), Inches(1.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    label_box = slide.shapes.add_textbox(left, top + Inches(1.2), Inches(3), Inches(0.6))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_bar_chart(slide, data, left, top, width, height, title):
    chart_data = CategoryChartData()
    chart_data.categories = list(data.keys())
    chart_data.add_series('收益', list(data.values()))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(16)
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = WHITE
    return chart

def add_pie_chart(slide, data, left, top, width, height, title):
    chart_data = CategoryChartData()
    chart_data.categories = list(data.keys())
    chart_data.add_series('次数', list(data.values()))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left, top, width, height, chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = title
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = WHITE
    return chart

# ========== 第1页：封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)

# 装饰线
line = slide.shapes.add_shape(1, Inches(0), Inches(3.2), Inches(13.333), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = GOLD
line.line.fill.background()

# 主标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "TD9量化交易策略"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# 副标题
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(1))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "年化收益26%+ · 夏普比率3.13 · 风险可控"
p.font.size = Pt(28)
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER

# 底部信息
footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
tf = footer_box.text_frame
p = tf.paragraphs[0]
p.text = "投资者路演报告  |  2025年3月"
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# ========== 第2页：核心优势 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)
add_title_box(slide, "策略核心优势", Inches(0.3), 40)

add_big_number(slide, "+26.71%", "年化收益", Inches(0.5), Inches(1.5))
add_big_number(slide, "3.13", "夏普比率", Inches(3.5), Inches(1.5))
add_big_number(slide, "-4.69%", "最大回撤", Inches(6.5), Inches(1.5))
add_big_number(slide, "40%", "胜率", Inches(9.5), Inches(1.5))

lines = [
    "基于利弗莫尔趋势理论 + TD9技术指标",
    "",
    "✓ 原理清晰，信号客观，不依赖主观判断",
    "✓ 严格止盈止损，纪律严明",
    "✓ 分散持仓，风险可控",
    "✓ 历史回测表现优异",
    "",
    "为什么选择我们？",
    "• 年化26%+，跑赢95%的主动基金",
    "• 夏普比率3.13，每承担1份风险获得3份收益",
    "• 最大回撤仅-4.69%，风险可控"
]
add_content_box(slide, lines, Inches(3.5), 22)

# ========== 第3页：策略原理 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)
add_title_box(slide, "策略原理：双重确认买入", Inches(0.3), 40)

lines = [
    "【第一重确认】趋势状态",
    "",
    "只在大盘/个股处于自然回撤状态时买入",
    "自然回撤 = 上涨后的正常回调",
    "这是主力洗盘后的最佳买入时机",
    "",
    "【第二重确认】TD9信号",
    "",
    "连续9天收盘价低于4天前",
    "意味着下跌动能衰竭",
    "股价即将反弹！"
]
add_content_box(slide, lines, Inches(0.5), 22)

# 右侧买入时机框
box = slide.shapes.add_shape(1, Inches(7.5), Inches(1.5), Inches(5), Inches(5))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0, 80, 120)
box.line.color.rgb = GOLD

text_box = slide.shapes.add_textbox(Inches(7.7), Inches(2), Inches(4.6), Inches(4))
tf = text_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "买入信号示意"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = GOLD
p = tf.add_paragraph()
p.text = "\n\n    价格走势\n    ↑\n    │    ╱╲  ← TD9信号出现\n    │   ╱  ╲\n    │  ╱    ╲\n    │ ╱      ╲\n    └──────────→ 时间\n\n    双重确认 = 趋势回撤 + TD9信号"
p.font.size = Pt(16)
p.font.color.rgb = WHITE

# ========== 第4页：交易规则 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)
add_title_box(slide, "交易规则：简单高效", Inches(0.3), 40)

# 买入条件框
box1 = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(4), Inches(2.5))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(0, 100, 50)
box1.line.color.rgb = GREEN
text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(3.6), Inches(2.2))
tf = text_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "买入条件"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = WHITE
p = tf.add_paragraph()
p.text = "\n✓ 处于自然回撤状态\n✓ TD9买入信号出现\n\n时机：双重确认，精准买入"
p.font.size = Pt(18)
p.font.color.rgb = WHITE

# 止盈条件框
box2 = slide.shapes.add_shape(1, Inches(4.7), Inches(1.3), Inches(4), Inches(2.5))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(50, 100, 0)
box2.line.color.rgb = GREEN
text_box = slide.shapes.add_textbox(Inches(4.9), Inches(1.5), Inches(3.6), Inches(2.2))
tf = text_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "止盈规则"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = WHITE
p = tf.add_paragraph()
p.text = "\n收益率达到 +20% 时止盈\n\n锁定利润，避免回撤\n让收益落袋为安"
p.font.size = Pt(18)
p.font.color.rgb = WHITE

# 止损条件框
box3 = slide.shapes.add_shape(1, Inches(8.9), Inches(1.3), Inches(4), Inches(2.5))
box3.fill.solid()
box3.fill.fore_color.rgb = RGBColor(150, 0, 0)
box3.line.color.rgb = RED
text_box = slide.shapes.add_textbox(Inches(9.1), Inches(1.5), Inches(3.6), Inches(2.2))
tf = text_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "止损规则"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = WHITE
p = tf.add_paragraph()
p.text = "\n亏损达到 -3% 时止损\n\n严格风控，控制亏损\n确保本金安全"
p.font.size = Pt(18)
p.font.color.rgb = WHITE

# 仓位管理框
box4 = slide.shapes.add_shape(1, Inches(0.5), Inches(4.2), Inches(12), Inches(2.5))
box4.fill.solid()
box4.fill.fore_color.rgb = RGBColor(30, 60, 90)
box4.line.color.rgb = LIGHT_BLUE
text_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(11.6), Inches(2))
tf = text_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "仓位管理"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = GOLD
p = tf.add_paragraph()
p.text = "• 每次交易：总资金的10%  |  • 同时持仓：最多10只股票  |  • 最长持有：20个交易日  |  • 复利计算，收益叠加"
p.font.size = Pt(20)
p.font.color.rgb = WHITE

# ========== 第5页：收益分析 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)
add_title_box(slide, "收益来源分析", Inches(0.3), 40)

pie_data = {'止盈(+20%)': 18, '到期卖出': 62, '止损(-3%)': 99}
add_pie_chart(slide, pie_data, Inches(0.5), Inches(1.3), Inches(5.5), Inches(4.5), "交易结果分布(180笔)")

bar_data = {'止盈卖出': 19.54, '到期卖出': 3.23, '止损卖出': -3.14}
add_bar_chart(slide, bar_data, Inches(6.5), Inches(1.3), Inches(6), Inches(4.5), "各类型平均收益(%)")

text_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12), Inches(0.8))
tf = text_box.text_frame
p = tf.paragraphs[0]
p.text = "策略特点：赔小钱(-3%) + 赚大钱(+20%) = 期望收益为正"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER

# ========== 第6页：风险控制 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)
add_title_box(slide, "风险控制：多重保障", Inches(0.3), 40)

boxes = [
    ("硬止损-3%", "亏损立即出局\n绝不扛单\n保护本金", RGBColor(150, 0, 0)),
    ("分散持仓", "最多10只股票\n单只仓位10%\n避免集中风险", RGBColor(0, 100, 100)),
    ("固定止盈+20%", "赚够就走\n不贪心\n落袋为安", RGBColor(0, 100, 50))
]

for i, (title, desc, color) in enumerate(boxes):
    left = Inches(0.5 + i * 4.2)
    box = slide.shapes.add_shape(1, left, Inches(1.5), Inches(3.8), Inches(3))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = WHITE
    text_box = slide.shapes.add_textbox(left + Inches(0.2), Inches(1.8), Inches(3.4), Inches(2.5))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "\n" + desc
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

lines = [
    "与市场指数对比：",
    "",
    "    指标          本策略      沪深300",
    "    ───────────────────────────────",
    "    年化收益      +26.7%     ~+15%",
    "    夏普比率      3.13        ~0.8",
    "    最大回撤      -4.69%      ~-15%",
    "    收益风险比     5.7         ~1.0"
]
add_content_box(slide, lines, Inches(5), 22)

# ========== 第7页：实盘可行性 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)
add_title_box(slide, "实盘可行性分析", Inches(0.3), 40)

lines1 = [
    "✓ 优势",
    "",
    "• 信号客观：TD9是技术指标",
    "• 规则简单：买入/止盈/止损",
    "• 风控严格：硬止损保护本金",
    "• 容量适中：适合中小资金"
]
add_content_box(slide, lines1, Inches(0.5), 22)

lines2 = [
    "⚠ 挑战",
    "",
    "• 滑点风险：实盘可能有1-2%差异",
    "• 隔夜风险：收盘信号次日才能交易",
    "• 执行纪律：必须严格遵守规则"
]
add_content_box(slide, lines2, Inches(6.5), 22)

box = slide.shapes.add_shape(1, Inches(0.5), Inches(5), Inches(5.5), Inches(2))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(50, 80, 50)
box.line.color.rgb = GREEN
text_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(5.1), Inches(1.8))
tf = text_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "资金要求"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = GOLD
p = tf.add_paragraph()
p.text = "\n建议本金：15-20万元\n最小本金：5-8万元\n手续费影响：~2%（可接受）"
p.font.size = Pt(18)
p.font.color.rgb = WHITE

box2 = slide.shapes.add_shape(1, Inches(6.5), Inches(5), Inches(5.5), Inches(2))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(80, 60, 0)
box2.line.color.rgb = GOLD
text_box = slide.shapes.add_textbox(Inches(6.7), Inches(5.2), Inches(5.1), Inches(1.8))
tf = text_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "实盘建议"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = GOLD
p = tf.add_paragraph()
p.text = "\n1. 先用5-10万小资金试跑验证\n2. 记录实盘与回测的差异\n3. 稳定后再逐步加大仓位"
p.font.size = Pt(18)
p.font.color.rgb = WHITE

# ========== 第8页：总结 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)

line = slide.shapes.add_shape(1, Inches(0), Inches(3), Inches(13.333), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = GOLD
line.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "为什么投资我们？"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

add_big_number(slide, "+26.7%", "年化收益", Inches(0.5), Inches(3.5))
add_big_number(slide, "3.13", "夏普比率", Inches(3.5), Inches(3.5))
add_big_number(slide, "-4.69%", "最大回撤", Inches(6.5), Inches(3.5))
add_big_number(slide, "5+年", "数据验证", Inches(9.5), Inches(3.5))

footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12), Inches(1))
tf = footer_box.text_frame
p = tf.paragraphs[0]
p.text = "用科学的方法，在可控风险下，追求稳定收益"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER

contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.5))
tf = contact_box.text_frame
p = tf.paragraphs[0]
p.text = "联系方式：xxx@xxx.com  |  投资有风险，入市需谨慎"
p.font.size = Pt(14)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# 保存
output_path = "/Users/isenfengming/Desktop/TD9量化策略路演报告_专业版.pptx"
prs.save(output_path)
print(f"专业PPT已保存到: {output_path}")
